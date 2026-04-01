import re
from pathlib import Path
from typing import Tuple


def extract_text(file_path: Path, mime_type: str) -> Tuple[str, int]:
    """
    Extract raw text and page count from a document.
    Returns (text, page_count).
    """
    suffix = file_path.suffix.lower()

    if suffix == ".pdf" or mime_type == "application/pdf":
        return _extract_pdf(file_path)
    elif suffix in (".docx",) or "wordprocessingml" in mime_type:
        return _extract_docx(file_path)
    elif suffix in (".pptx",) or "presentationml" in mime_type:
        return _extract_pptx(file_path)
    elif suffix in (".txt", ".md", ".csv") or mime_type.startswith("text/"):
        return _extract_text(file_path)
    else:
        # Best-effort: try as text
        try:
            return _extract_text(file_path)
        except Exception:
            return "", 0


def _extract_pdf(path: Path) -> Tuple[str, int]:
    import fitz  # PyMuPDF
    doc = fitz.open(str(path))
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages), len(pages)


def _extract_docx(path: Path) -> Tuple[str, int]:
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    text = "\n".join(paragraphs)
    # Estimate pages: ~400 words per page
    word_count = len(text.split())
    pages = max(1, word_count // 400)
    return text, pages


def _extract_pptx(path: Path) -> Tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides_text = []
    for slide in prs.slides:
        slide_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())
        if slide_parts:
            slides_text.append("\n".join(slide_parts))
    return "\n\n".join(slides_text), len(prs.slides)


def _extract_text(path: Path) -> Tuple[str, int]:
    text = path.read_text(encoding="utf-8", errors="replace")
    lines = len(text.splitlines())
    pages = max(1, lines // 50)
    return text, pages


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 100):
    """
    Split text into overlapping chunks for vector indexing.
    Tries to split on sentence boundaries.
    """
    # Clean up whitespace
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end >= len(text):
            chunks.append(text[start:])
            break

        # Try to find a sentence boundary near the end
        boundary = text.rfind(". ", start, end)
        if boundary == -1 or boundary < start + chunk_size // 2:
            boundary = text.rfind("\n", start, end)
        if boundary == -1 or boundary < start + chunk_size // 2:
            boundary = end

        chunks.append(text[start:boundary + 1].strip())
        start = boundary + 1 - overlap

    return [c for c in chunks if c.strip()]
